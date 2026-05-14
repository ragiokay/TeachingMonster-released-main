#!/usr/bin/env python3
"""
Slide Generator Module

Converts Layout JSON specifications into PowerPoint (PPTX) files.
Supports various slide layouts: TITLE, CONTENT, SECTION, TWO_CONTENT, COMPARISON, CUSTOM.
"""

import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

from . import text_utils


def hex_to_rgb(hex_color):
    """Convert hex string to RGBColor."""
    if not hex_color:
        return None
    hex_color = hex_color.lstrip('#')
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16)
    )


def apply_background(slide, color_hex):
    """Apply solid background color to slide."""
    if not color_hex:
        return
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(color_hex)


def apply_text_style(shape, text, style=None):
    """Apply text and simple styling to a shape."""
    if not style:
        style = {}

    if not hasattr(shape, "text_frame"):
        return

    tf = shape.text_frame
    tf.clear()

    # Text Auto-Scaling Logic
    if style.get('auto_size', True):
        # Disable PPT native auto-size to prevent inconsistencies
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.word_wrap = True

        # Calculate optimal font size using TextFitter
        if hasattr(shape, 'width') and hasattr(shape, 'height'):
            w_px = shape.width / 9525  # EMU to pixels
            h_px = shape.height / 9525

            start_sz = style.get('font_size', 32)
            fitter = text_utils.TextFitter()

            text_content = text
            if isinstance(text, list):
                text_content = "\n".join(text)

            optimal_size = fitter.fit_text(
                text_content,
                max_width_px=w_px,
                max_height_px=h_px,
                start_size=start_sz
            )

            style['font_size'] = optimal_size

    p = tf.paragraphs[0]
    p.text = text

    font = p.font
    if style.get('font_size'):
        font.size = Pt(style['font_size'])
    if style.get('bold'):
        font.bold = True
    if style.get('italic'):
        font.italic = True
    if style.get('color'):
        font.color.rgb = hex_to_rgb(style['color'])

    align_map = {
        'LEFT': PP_ALIGN.LEFT,
        'CENTER': PP_ALIGN.CENTER,
        'RIGHT': PP_ALIGN.RIGHT
    }
    if style.get('alignment'):
        p.alignment = align_map.get(style['alignment'], PP_ALIGN.LEFT)

    if style.get('link'):
        for run in p.runs:
            run.hyperlink.address = style['link']


def enforce_slide_layout(slide, prs, style=None):
    """
    Force standard placeholders to fit the actual slide aspect ratio.
    """
    if style is None:
        style = {}
    metrics = style.get("layout_metrics", {})

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Defaults
    DEFAULT_MARGIN_RATIO = 0.075
    DEFAULT_TOP_RATIO = 0.05
    DEFAULT_TITLE_H_RATIO = 0.15
    DEFAULT_GAP_INCH = 0.5

    margin_ratio = metrics.get("margin_h_ratio", DEFAULT_MARGIN_RATIO)
    top_ratio = metrics.get("margin_top_title_ratio", DEFAULT_TOP_RATIO)
    title_h_ratio = metrics.get("title_height_ratio", DEFAULT_TITLE_H_RATIO)
    gap_inch = metrics.get("gap_inch", DEFAULT_GAP_INCH)

    MARGIN_H = int(slide_width * margin_ratio)
    MARGIN_TOP_TITLE = int(slide_height * top_ratio)
    TITLE_HEIGHT = int(slide_height * title_h_ratio)

    CONTENT_GAP = int(slide_height * 0.05)
    CONTENT_TOP = MARGIN_TOP_TITLE + TITLE_HEIGHT + CONTENT_GAP
    CONTENT_HEIGHT = slide_height - CONTENT_TOP - int(slide_height * 0.05)
    CONTENT_WIDTH = slide_width - (MARGIN_H * 2)

    ph_indices = [p.placeholder_format.idx for p in slide.placeholders]

    # TITLE (idx 0)
    if 0 in ph_indices:
        title = slide.placeholders[0]
        title.left = MARGIN_H
        title.width = CONTENT_WIDTH
        title.top = MARGIN_TOP_TITLE
        title.height = TITLE_HEIGHT

        # Title Slide special handling
        if len(ph_indices) == 2 and 1 in ph_indices and slide.slide_layout == prs.slide_layouts[0]:
            title_height_large = int(slide_height * 0.25)
            title.top = int(slide_height * 0.3)
            title.height = title_height_large

            sub = slide.placeholders[1]
            sub.left = MARGIN_H
            sub.width = CONTENT_WIDTH
            sub.top = title.top + title_height_large + int(slide_height * 0.02)
            sub.height = int(slide_height * 0.15)
            return

    GAP_EMU = int(gap_inch * 914400)
    HALF_WIDTH = int((CONTENT_WIDTH - GAP_EMU) / 2)

    def fix_content_ph(ph, x, w, y=CONTENT_TOP, h=CONTENT_HEIGHT):
        ph.left = x
        ph.width = w
        ph.top = y
        ph.height = h

    # CONTENT SLIDE (Layout 1)
    if len(ph_indices) == 2 and 1 in ph_indices and slide.slide_layout == prs.slide_layouts[1]:
        fix_content_ph(slide.placeholders[1], MARGIN_H, CONTENT_WIDTH)

    # SECTION HEADER (Layout 2)
    elif len(ph_indices) >= 2 and slide.slide_layout == prs.slide_layouts[2]:
        t_h = int(slide_height * 0.2)
        slide.placeholders[0].top = int(slide_height * 0.4) - t_h
        slide.placeholders[0].height = t_h

        if 1 in ph_indices:
            p1 = slide.placeholders[1]
            p1.left = MARGIN_H
            p1.width = CONTENT_WIDTH
            p1.top = slide.placeholders[0].top + t_h + int(slide_height * 0.05)
            p1.height = int(slide_height * 0.2)

    # TWO CONTENT (Layout 3)
    elif len(ph_indices) >= 3 and slide.slide_layout == prs.slide_layouts[3]:
        if 1 in ph_indices:
            fix_content_ph(slide.placeholders[1], MARGIN_H, HALF_WIDTH)
        if 2 in ph_indices:
            fix_content_ph(slide.placeholders[2], MARGIN_H + HALF_WIDTH + GAP_EMU, HALF_WIDTH)

    # COMPARISON (Layout 4)
    elif len(ph_indices) >= 5 and slide.slide_layout == prs.slide_layouts[4]:
        HEAD_TOP = CONTENT_TOP
        HEAD_HEIGHT = int(slide_height * 0.08)
        BODY_TOP = HEAD_TOP + HEAD_HEIGHT
        BODY_HEIGHT = CONTENT_HEIGHT - HEAD_HEIGHT

        if 1 in ph_indices:
            fix_content_ph(slide.placeholders[1], MARGIN_H, HALF_WIDTH, HEAD_TOP, HEAD_HEIGHT)
        if 2 in ph_indices:
            fix_content_ph(slide.placeholders[2], MARGIN_H, HALF_WIDTH, BODY_TOP, BODY_HEIGHT)
        if 3 in ph_indices:
            fix_content_ph(slide.placeholders[3], MARGIN_H + HALF_WIDTH + GAP_EMU, HALF_WIDTH, HEAD_TOP, HEAD_HEIGHT)
        if 4 in ph_indices:
            fix_content_ph(slide.placeholders[4], MARGIN_H + HALF_WIDTH + GAP_EMU, HALF_WIDTH, BODY_TOP, BODY_HEIGHT)


def create_title_slide(prs, content, style=None):
    """Standard Title Slide."""
    if style is None:
        style = {}
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    enforce_slide_layout(slide, prs, style)
    apply_background(slide, style.get("background"))

    apply_text_style(slide.shapes.title, content.get("title", ""), style.get("title_style"))
    apply_text_style(slide.placeholders[1], content.get("subtitle", ""), style.get("subtitle_style"))
    return slide


def create_content_slide(prs, content, style=None):
    """Standard Title + Content Slide."""
    if style is None:
        style = {}
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    enforce_slide_layout(slide, prs, style)
    apply_background(slide, style.get("background"))

    apply_text_style(slide.shapes.title, content.get("title", ""), style.get("title_style"))

    tf = slide.placeholders[1].text_frame
    tf.clear()
    points = content.get("body", [])
    if isinstance(points, str):
        points = [points]

    align_map = {'LEFT': PP_ALIGN.LEFT, 'CENTER': PP_ALIGN.CENTER, 'RIGHT': PP_ALIGN.RIGHT}
    body_style = style.get("body_style", {})
    for i, point in enumerate(points):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = point
        if body_style.get("font_size"):
            p.font.size = Pt(body_style["font_size"])
        if body_style.get("color"):
            p.font.color.rgb = hex_to_rgb(body_style["color"])
        if body_style.get("alignment"):
            p.alignment = align_map.get(body_style["alignment"], PP_ALIGN.CENTER)

    return slide


def create_section_slide(prs, content, style=None):
    """Section Header Slide."""
    if style is None:
        style = {}
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    enforce_slide_layout(slide, prs, style)
    apply_background(slide, style.get("background"))

    apply_text_style(slide.shapes.title, content.get("title", ""), style.get("title_style"))
    if len(slide.placeholders) > 1:
        apply_text_style(slide.placeholders[1], content.get("text", ""), style.get("text_style"))
    return slide


def create_custom_slide(prs, content, style=None):
    """
    Freeform Slide (Blank Layout).
    Elements are placed via (x, y, w, h) in inches.
    """
    if style is None:
        style = {}
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 6 = Blank
    apply_background(slide, style.get("background"))

    elements = content.get("elements", [])
    for elem in elements:
        x = Inches(elem.get("x", 0))
        y = Inches(elem.get("y", 0))
        w = Inches(elem.get("w", 1))
        h = Inches(elem.get("h", 1))

        elem_type = elem.get("type", "textbox")

        if elem_type == "textbox":
            textbox = slide.shapes.add_textbox(x, y, w, h)
            apply_text_style(textbox, elem.get("text", ""), elem.get("style"))

        elif elem_type == "shape":
            shape_type_map = {"RECTANGLE": MSO_SHAPE.RECTANGLE, "OVAL": MSO_SHAPE.OVAL}
            st = shape_type_map.get(elem.get("shape_type"), MSO_SHAPE.RECTANGLE)
            shape = slide.shapes.add_shape(st, x, y, w, h)

            if elem.get("color"):
                shape.fill.solid()
                shape.fill.fore_color.rgb = hex_to_rgb(elem["color"])
                shape.line.fill.background()

            if elem.get("text"):
                apply_text_style(shape, elem["text"], elem.get("style"))

        elif elem_type == "image":
            image_path = elem.get("path")
            prompt = elem.get("prompt")

            if not image_path and prompt:
                print(f"Warning: Image path missing for prompt: '{prompt}'.")

            if image_path and os.path.exists(image_path):
                slide.shapes.add_picture(image_path, x, y, w, h)
            else:
                print(f"Warning: Image path not found: {image_path}")

        elif elem_type == "math":
            latex = elem.get("latex", "")
            elem_style = elem.get("style", {})
            font_size = elem_style.get("font_size", 24)
            color = elem_style.get("color", "#000000")

            if latex:
                from . import math_renderer
                renderer = math_renderer.MathRenderer()

                math_image_path = renderer.render_latex(
                    latex=latex,
                    font_size=font_size,
                    color=color,
                    dpi=300
                )

                if math_image_path and os.path.exists(math_image_path):
                    slide.shapes.add_picture(math_image_path, x, y, w, h)
                else:
                    print(f"Warning: Math rendering failed for: {latex}")
                    box = slide.shapes.add_textbox(x, y, w, h)
                    tf = box.text_frame
                    tf.text = f"[Math: {latex}]"
            else:
                print(f"Warning: Math element missing 'latex' field")

        elif elem_type == "table":
            rows = elem.get("rows", 2)
            cols = elem.get("cols", 2)
            table_shape = slide.shapes.add_table(rows, cols, x, y, w, h)
            table = table_shape.table

            data = elem.get("data", [])
            for r in range(min(rows, len(data))):
                for c in range(min(cols, len(data[r]))):
                    cell = table.cell(r, c)
                    cell.text = str(data[r][c])
                    if elem.get("style"):
                        apply_text_style(cell, str(data[r][c]), elem.get("style"))

        elif elem_type == "chart":
            chart_type_map = {
                "BAR_CLUSTERED": XL_CHART_TYPE.BAR_CLUSTERED,
                "COLUMN_CLUSTERED": XL_CHART_TYPE.COLUMN_CLUSTERED,
                "LINE": XL_CHART_TYPE.LINE,
                "PIE": XL_CHART_TYPE.PIE
            }
            c_type = chart_type_map.get(elem.get("chart_type"), XL_CHART_TYPE.COLUMN_CLUSTERED)

            chart_data = CategoryChartData()
            categories = elem.get("categories", [])
            chart_data.categories = categories

            series_list = elem.get("series", [])
            for series in series_list:
                chart_data.add_series(series.get("name", "Series"), series.get("values", []))

            chart = slide.shapes.add_chart(
                c_type, x, y, w, h, chart_data
            ).chart

            if elem.get("legend", True):
                chart.has_legend = True
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM

    return slide


def create_two_content_slide(prs, content, style=None):
    """
    Two Content Slide (Layout 3).
    Placeholders: 0=Title, 1=Left Content, 2=Right Content
    """
    if style is None:
        style = {}
    slide = prs.slides.add_slide(prs.slide_layouts[3])
    enforce_slide_layout(slide, prs, style)
    apply_background(slide, style.get("background"))

    apply_text_style(slide.shapes.title, content.get("title", ""), style.get("title_style"))

    tf_left = slide.placeholders[1].text_frame
    tf_left.clear()
    left_points = content.get("left", [])
    if isinstance(left_points, str):
        left_points = [left_points]

    align_map = {'LEFT': PP_ALIGN.LEFT, 'CENTER': PP_ALIGN.CENTER, 'RIGHT': PP_ALIGN.RIGHT}
    body_style = style.get("body_style", {})
    for i, point in enumerate(left_points):
        p = tf_left.add_paragraph() if i > 0 else tf_left.paragraphs[0]
        p.text = point
        if body_style.get("font_size"):
            p.font.size = Pt(body_style["font_size"])
        if body_style.get("color"):
            p.font.color.rgb = hex_to_rgb(body_style["color"])
        if body_style.get("alignment"):
            p.alignment = align_map.get(body_style["alignment"], PP_ALIGN.LEFT)

    tf_right = slide.placeholders[2].text_frame
    tf_right.clear()
    right_points = content.get("right", [])
    if isinstance(right_points, str):
        right_points = [right_points]

    for i, point in enumerate(right_points):
        p = tf_right.add_paragraph() if i > 0 else tf_right.paragraphs[0]
        p.text = point
        if body_style.get("font_size"):
            p.font.size = Pt(body_style["font_size"])
        if body_style.get("color"):
            p.font.color.rgb = hex_to_rgb(body_style["color"])
        if body_style.get("alignment"):
            p.alignment = align_map.get(body_style["alignment"], PP_ALIGN.LEFT)

    return slide


def create_comparison_slide(prs, content, style=None):
    """
    Comparison Slide (Layout 4).
    Placeholders: 0=Title, 1=Left Header, 2=Left Content, 3=Right Header, 4=Right Content
    """
    if style is None:
        style = {}
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    enforce_slide_layout(slide, prs, style)
    apply_background(slide, style.get("background"))

    apply_text_style(slide.shapes.title, content.get("title", ""), style.get("title_style"))

    text_style = style.get("text_style", {})
    apply_text_style(slide.placeholders[1], content.get("left_header", ""), text_style)
    apply_text_style(slide.placeholders[3], content.get("right_header", ""), text_style)

    align_map = {'LEFT': PP_ALIGN.LEFT, 'CENTER': PP_ALIGN.CENTER, 'RIGHT': PP_ALIGN.RIGHT}
    body_style = style.get("body_style", {})

    tf_left = slide.placeholders[2].text_frame
    tf_left.clear()
    left_points = content.get("left", [])
    if isinstance(left_points, str):
        left_points = [left_points]

    for i, point in enumerate(left_points):
        p = tf_left.add_paragraph() if i > 0 else tf_left.paragraphs[0]
        p.text = point
        if body_style.get("font_size"):
            p.font.size = Pt(body_style["font_size"])
        if body_style.get("color"):
            p.font.color.rgb = hex_to_rgb(body_style["color"])
        if body_style.get("alignment"):
            p.alignment = align_map.get(body_style["alignment"], PP_ALIGN.LEFT)

    tf_right = slide.placeholders[4].text_frame
    tf_right.clear()
    right_points = content.get("right", [])
    if isinstance(right_points, str):
        right_points = [right_points]

    for i, point in enumerate(right_points):
        p = tf_right.add_paragraph() if i > 0 else tf_right.paragraphs[0]
        p.text = point
        if body_style.get("font_size"):
            p.font.size = Pt(body_style["font_size"])
        if body_style.get("color"):
            p.font.color.rgb = hex_to_rgb(body_style["color"])
        if body_style.get("alignment"):
            p.alignment = align_map.get(body_style["alignment"], PP_ALIGN.LEFT)

    return slide


def generate_slides(json_input_path: str, output_dir: str):
    """
    Generate PPTX files from a JSON specification file.

    Args:
        json_input_path: Path to JSON file with slide specifications
        output_dir: Directory to save generated PPTX files
    """
    with open(json_input_path, 'r') as f:
        slides_data = json.load(f)

    os.makedirs(output_dir, exist_ok=True)

    for i, slide_data in enumerate(slides_data):
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        slide_id = slide_data.get("id", f"slide_{i}")
        layout = slide_data.get("layout", "CONTENT").upper()
        content = slide_data.get("content", {})
        style = slide_data.get("style", {})

        # If content has 'elements' list, force CUSTOM render logic
        if "elements" in content and isinstance(content["elements"], list) and len(content["elements"]) > 0:
            create_custom_slide(prs, content, style)
        elif layout == "TITLE":
            create_title_slide(prs, content, style)
        elif layout == "SECTION":
            create_section_slide(prs, content, style)
        elif layout == "TWO_CONTENT":
            create_two_content_slide(prs, content, style)
        elif layout == "COMPARISON":
            create_comparison_slide(prs, content, style)
        elif layout == "CUSTOM":
            create_custom_slide(prs, content, style)
        else:
            create_content_slide(prs, content, style)

        output_path = os.path.join(output_dir, f"{slide_id}.pptx")
        prs.save(output_path)
        print(f"Generated {output_path}")


def generate_single_slide(slide_data: dict, output_dir: str) -> str:
    """
    Generate a single PPTX file from a slide specification dict.

    Args:
        slide_data: Dict containing slide layout, content, and style
        output_dir: Directory to save the generated PPTX file

    Returns:
        Path to the generated PPTX file
    """
    os.makedirs(output_dir, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_id = slide_data.get("id", "slide_0")
    layout = slide_data.get("layout", "CONTENT").upper()
    content = slide_data.get("content", {})
    # Normalise: LLM sometimes returns content as a bare list of elements
    if isinstance(content, list):
        content = {"elements": content}
    style = slide_data.get("style", {})

    if "elements" in content and isinstance(content["elements"], list) and len(content["elements"]) > 0:
        create_custom_slide(prs, content, style)
    elif layout == "TITLE":
        create_title_slide(prs, content, style)
    elif layout == "SECTION":
        create_section_slide(prs, content, style)
    elif layout == "TWO_CONTENT":
        create_two_content_slide(prs, content, style)
    elif layout == "COMPARISON":
        create_comparison_slide(prs, content, style)
    elif layout == "CUSTOM":
        create_custom_slide(prs, content, style)
    else:
        create_content_slide(prs, content, style)

    output_path = os.path.join(output_dir, f"{slide_id}.pptx")
    prs.save(output_path)

    return output_path
